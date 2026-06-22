#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генерация презентации проекта "Спадчына" (CultCode).
Цветовая палитра взята из проекта (index.css).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Цветовая палитра проекта ---
COLORS = {
    "primary": RGBColor(15, 81, 50),        # #0f5132
    "primary_light": RGBColor(230, 242, 236),
    "accent": RGBColor(185, 28, 28),        # #b91c1c
    "accent_light": RGBColor(253, 232, 232),
    "bg": RGBColor(248, 245, 239),          # #f8f5ef
    "surface": RGBColor(255, 255, 255),     # #ffffff
    "text": RGBColor(31, 41, 51),           # #1f2933
    "muted": RGBColor(107, 114, 128),       # приглушённый текст
    "border": RGBColor(216, 218, 221),
}

# --- Настройки ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.6)


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


def add_background(slide, color):
    """Добавляет цветной фон слайда."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        int(SLIDE_WIDTH), int(SLIDE_HEIGHT)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Перемещаем на задний план
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_title(slide, text, left, top, width, height,
              font_size=44, bold=True, color=None, align=PP_ALIGN.LEFT):
    """Добавляет заголовок."""
    if color is None:
        color = COLORS["text"]
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    return box


def add_body(slide, text, left, top, width, height,
             font_size=18, color=None, align=PP_ALIGN.LEFT, line_spacing=1.3):
    """Добавляет основной текст."""
    if color is None:
        color = COLORS["text"]
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    p.line_spacing = line_spacing
    return box


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=18, color=None, bullet_color=None):
    """Добавляет маркированный список."""
    if color is None:
        color = COLORS["text"]
    if bullet_color is None:
        bullet_color = COLORS["primary"]
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.line_spacing = 1.4
        p.space_after = Pt(8)
        # Маркер
        p.font.name = "Arial"
    return box


def add_card(slide, left, top, width, height, color, icon_text, title, description, text_color=None):
    """Добавляет карточку с иконкой, заголовком и описанием."""
    if text_color is None:
        text_color = COLORS["text"]
    # Фон карточки
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLORS["border"]
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05

    # Иконка/буква
    icon_size = Inches(0.55)
    icon_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(left + Inches(0.2)), int(top + Inches(0.2)),
        int(icon_size), int(icon_size)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = COLORS["primary"]
    icon_box.line.fill.background()
    tf = icon_box.text_frame
    tf.paragraphs[0].text = icon_text
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["surface"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Arial"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Заголовок
    add_title(slide, title,
              left + Inches(0.85), top + Inches(0.18),
              width - Inches(1.05), Inches(0.35),
              font_size=16, bold=True, color=text_color)

    # Описание
    add_body(slide, description,
             left + Inches(0.85), top + Inches(0.55),
             width - Inches(1.05), height - Inches(0.75),
             font_size=13, color=COLORS["muted"])


def add_decorative_strip(slide, color=None, top=True):
    """Декоративная полоса сверху или снизу."""
    if color is None:
        color = COLORS["primary"]
    height = Inches(0.18)
    if top:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            int(SLIDE_WIDTH), int(height)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, int(SLIDE_HEIGHT - height),
            int(SLIDE_WIDTH), int(height)
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ============================================================
    # Слайд 1: Титульный
    # ============================================================
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])

    # Декоративные элементы
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(Inches(8.5)), 0,
        int(Inches(4.833)), int(SLIDE_HEIGHT)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)

    add_title(slide, "Спадчына",
              MARGIN, Inches(2.2),
              Inches(7.5), Inches(1),
              font_size=64, bold=True, color=COLORS["primary"])

    add_title(slide, "Интерактивный гид по культурному наследию Беларуси",
              MARGIN, Inches(3.3),
              Inches(7.5), Inches(1),
              font_size=24, bold=False, color=COLORS["text"])

    add_body(slide, "Образовательная платформа с викторинами, дуэлями, достижениями и чатом",
             MARGIN, Inches(4.3),
             Inches(7.5), Inches(0.8),
             font_size=18, color=COLORS["muted"])

    add_body(slide, "CultCode / Repablik",
             MARGIN, Inches(6.4),
             Inches(7.5), Inches(0.5),
             font_size=14, color=COLORS["muted"])

    # ============================================================
    # Слайд 2: Что это?
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["primary"], top=True)

    add_title(slide, "Что такое «Спадчына»?",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    add_body(slide,
             "Это веб-платформа, где школьники и молодёжь изучают культурное наследие Беларуси не по учебнику, а через интерактив: статьи, тесты, соревнования и общение.",
             MARGIN, Inches(1.5),
             Inches(12), Inches(1),
             font_size=20, color=COLORS["text"])

    cards = [
        ("📚", "Библиотека знаний", "50+ статей о достопримечательностях, истории, культуре и природе"),
        ("🎯", "Викторины", "Проверка знаний после каждой статьи с начислением баллов"),
        ("🏆", "Геймификация", "Достижения, таблица лидеров, магазин декораций"),
        ("⚔️", "Дуэли и команды", "PvP-викторины и групповые битвы по 6-значному коду"),
        ("💬", "Чат", "Личные сообщения с голосовым вводом на трёх языках"),
        ("🛡️", "Админ-панель", "Управление контентом и пользователями"),
    ]

    card_width = Inches(3.7)
    card_height = Inches(1.6)
    start_x = Inches(0.5)
    start_y = Inches(2.7)
    gap_x = Inches(0.25)
    gap_y = Inches(0.3)

    for i, (icon, title, desc) in enumerate(cards):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        add_card(slide, x, y, card_width, card_height,
                 COLORS["surface"], icon, title, desc)

    # ============================================================
    # Слайд 3: Проблема и решение
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["accent"], top=True)

    add_title(slide, "Проблема и решение",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    # Левая колонка — проблема
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(MARGIN), int(Inches(1.6)),
        int(Inches(5.8)), int(Inches(4.8))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent_light"]
    shape.line.color.rgb = COLORS["accent"]
    shape.line.width = Pt(2)
    shape.adjustments[0] = 0.05

    add_title(slide, "Проблема",
              MARGIN + Inches(0.3), Inches(1.85),
              Inches(5.2), Inches(0.5),
              font_size=26, bold=True, color=COLORS["accent"])

    problems = [
        "• Традиционные уроки быстро утрачивают внимание",
        "• Мало интерактивных ресурсов по белорусской культуре",
        "• Сложно мотивировать к самостоятельному изучению",
        "• Отсутствует соревновательный элемент",
    ]
    for idx, txt in enumerate(problems):
        add_body(slide, txt,
                 MARGIN + Inches(0.3), Inches(2.5) + idx * Inches(0.55),
                 Inches(5.2), Inches(0.5),
                 font_size=18, color=COLORS["text"])

    # Правая колонка — решение
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(Inches(7.0)), int(Inches(1.6)),
        int(Inches(5.8)), int(Inches(4.8))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary_light"]
    shape.line.color.rgb = COLORS["primary"]
    shape.line.width = Pt(2)
    shape.adjustments[0] = 0.05

    add_title(slide, "Решение",
              Inches(7.3), Inches(1.85),
              Inches(5.2), Inches(0.5),
              font_size=26, bold=True, color=COLORS["primary"])

    solutions = [
        "• Игровое обучение: баллы, достижения, лидерборд",
        "• Понятный контент с категориями и сложностью",
        "• Социальные механики: дуэли, команды, чат",
        "• Голосовой ввод и доступность с любого устройства",
    ]
    for idx, txt in enumerate(solutions):
        add_body(slide, txt,
                 Inches(7.3), Inches(2.5) + idx * Inches(0.55),
                 Inches(5.2), Inches(0.5),
                 font_size=18, color=COLORS["text"])

    # ============================================================
    # Слайд 4: Стек технологий
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["primary"], top=True)

    add_title(slide, "Технологический стек",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    # Frontend card
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(MARGIN), int(Inches(1.6)),
        int(Inches(5.8)), int(Inches(5.0))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["surface"]
    shape.line.color.rgb = COLORS["primary"]
    shape.line.width = Pt(2)
    shape.adjustments[0] = 0.05

    add_title(slide, "Фронтенд",
              MARGIN + Inches(0.3), Inches(1.85),
              Inches(5.2), Inches(0.5),
              font_size=28, bold=True, color=COLORS["primary"])

    frontend_items = [
        "React 18 — UI-компоненты",
        "Vite 5 — сборка и dev-сервер",
        "Чистый CSS + CSS Variables",
        "react-speech-recognition",
        "FontAwesome 7 — иконки",
        "ESLint — качество кода",
    ]
    for idx, txt in enumerate(frontend_items):
        add_body(slide, "• " + txt,
                 MARGIN + Inches(0.3), Inches(2.55) + idx * Inches(0.55),
                 Inches(5.2), Inches(0.5),
                 font_size=18, color=COLORS["text"])

    # Backend card
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(Inches(7.0)), int(Inches(1.6)),
        int(Inches(5.8)), int(Inches(5.0))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["surface"]
    shape.line.color.rgb = COLORS["accent"]
    shape.line.width = Pt(2)
    shape.adjustments[0] = 0.05

    add_title(slide, "Бэкенд",
              Inches(7.3), Inches(1.85),
              Inches(5.2), Inches(0.5),
              font_size=28, bold=True, color=COLORS["accent"])

    backend_items = [
        "Go 1.26 — HTTP API",
        "SQLite — файловая БД",
        "JWT — аутентификация",
        "bcrypt — хэширование",
        "Telegram Bot API",
        "CORS + REST",
    ]
    for idx, txt in enumerate(backend_items):
        add_body(slide, "• " + txt,
                 Inches(7.3), Inches(2.55) + idx * Inches(0.55),
                 Inches(5.2), Inches(0.5),
                 font_size=18, color=COLORS["text"])

    # ============================================================
    # Слайд 5: Архитектура
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["primary"], top=True)

    add_title(slide, "Архитектура системы",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    # Схема слоёв
    layers = [
        ("Пользователь", "Браузер / мобильное устройство", COLORS["surface"], Inches(0.8)),
        ("Фронтенд", "React + Vite\nhttp://localhost:5173", COLORS["primary_light"], Inches(2.1)),
        ("API", "Go REST API\n/api/* → http://127.0.0.1:8081", COLORS["primary"], Inches(3.4)),
        ("База данных", "SQLite\ncultcode.db", COLORS["accent"], Inches(4.7)),
    ]

    for title, desc, color, top in layers:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(Inches(3.5)), int(top),
            int(Inches(6.3)), int(Inches(1.0))
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.08

        text_color = COLORS["surface"] if color in (COLORS["primary"], COLORS["accent"]) else COLORS["text"]
        add_title(slide, title,
                  Inches(3.7), top + Inches(0.08),
                  Inches(5.9), Inches(0.35),
                  font_size=20, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        add_body(slide, desc,
                 Inches(3.7), top + Inches(0.42),
                 Inches(5.9), Inches(0.5),
                 font_size=14, color=text_color, align=PP_ALIGN.CENTER)

    # Стрелки между слоями
    for i in range(3):
        top = Inches(1.8) + i * Inches(1.3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            int(Inches(6.25)), int(top),
            int(Inches(0.6)), int(Inches(0.35))
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["muted"]
        shape.line.fill.background()

    add_body(slide, "Vite проксирует запросы /api на Go-сервер. Всё хранится в одной SQLite-базе.",
             MARGIN, Inches(6.0),
             Inches(12), Inches(0.8),
             font_size=16, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    # ============================================================
    # Слайд 6: Ключевые возможности (1)
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["accent"], top=True)

    add_title(slide, "Ключевые возможности",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    feature_cards = [
        ("🗺️", "Карточки мест", "Статьи с фото, категорией, сложностью и временем чтения"),
        ("❓", "Викторины", "Single/multiple choice, подсчёт результата, бонусы за идеал"),
        ("📊", "Прогресс", "Круговая диаграмма, список пройденных мест, достижения"),
        ("🏅", "Достижения", "12 значков: от «Первые шаги» до «Легенда»"),
        ("🥇", "Лидерборд", "Топ-50 по правильным ответам, профили участников"),
        ("🛍️", "Магазин", "Декорации за баллы: карточки, рамки, стикеры"),
    ]

    card_width = Inches(3.7)
    card_height = Inches(1.6)
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    gap_x = Inches(0.25)
    gap_y = Inches(0.3)

    for i, (icon, title, desc) in enumerate(feature_cards):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        add_card(slide, x, y, card_width, card_height,
                 COLORS["surface"], icon, title, desc)

    # ============================================================
    # Слайд 7: Ключевые возможности (2) — социальные
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["primary"], top=True)

    add_title(slide, "Социальные механики",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    social_cards = [
        ("⚔️", "Дуэли 1×1", "10 вопросов на время, бонус за скорость. Победитель получает +40 баллов и редкую рамку"),
        ("👥", "Командные битвы", "Комната по 6-значному коду. Категория + 10–20 вопросов. Топ-3: 60/40/20 баллов"),
        ("💬", "Чат", "Личные сообщения, поиск диалогов, непрочитанные сообщения, голосовой ввод RU/BY/EN"),
        ("🤖", "Telegram-бот", "Приём предложений от пользователей прямо в Telegram"),
        ("🛡️", "Админ-панель", "Создание статей из plain text, управление пользователями"),
        ("💾", "Офлайн-фолбэк", "50+ статей зашиты локально — работает даже без сервера"),
    ]

    for i, (icon, title, desc) in enumerate(social_cards):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        add_card(slide, x, y, card_width, card_height,
                 COLORS["surface"], icon, title, desc)

    # ============================================================
    # Слайд 8: Геймификация
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["accent"], top=True)

    add_title(slide, "Геймификация и баллы",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    scoring = [
        ("Викторина", "правильных × 10 + 20 за идеал / +10 за ≥70%"),
        ("Дуэль", "100 + остаток_времени × 5 за вопрос"),
        ("Победа в дуэли", "+40 баллов"),
        ("Командная битва", "100 за правильный ответ"),
        ("Топ-3 в битве", "+60 / +40 / +20 баллов"),
    ]

    for idx, (name, rule) in enumerate(scoring):
        y = Inches(1.6) + idx * Inches(0.95)

        # Название
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(MARGIN), int(y),
            int(Inches(3.0)), int(Inches(0.75))
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["primary"]
        shape.line.fill.background()
        shape.adjustments[0] = 0.1
        add_body(slide, name,
                 MARGIN + Inches(0.15), y + Inches(0.15),
                 Inches(2.7), Inches(0.45),
                 font_size=16, color=COLORS["surface"], align=PP_ALIGN.CENTER)

        # Правило
        add_body(slide, rule,
                 Inches(4.0), y + Inches(0.15),
                 Inches(8.5), Inches(0.6),
                 font_size=18, color=COLORS["text"])

    add_title(slide, "12 достижений",
              MARGIN, Inches(6.1),
              Inches(12), Inches(0.5),
              font_size=22, bold=True, color=COLORS["primary"])

    achievements_text = "Первые шаги • Искатель • Путешественник • Исследователь • Внимательный • Эрудит • Эксперт • Безупречно • Мастер идеала • Новичок • Профи • Легенда"
    add_body(slide, achievements_text,
             MARGIN, Inches(6.55),
             Inches(12), Inches(0.6),
             font_size=14, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    # ============================================================
    # Слайд 9: Цветовая палитра
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["primary"], top=True)

    add_title(slide, "Цветовая палитра",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    palette = [
        ("Основной зелёный", "#0f5132", COLORS["primary"], COLORS["surface"]),
        ("Акцентный красный", "#b91c1c", COLORS["accent"], COLORS["surface"]),
        ("Фон страницы", "#f8f5ef", COLORS["bg"], COLORS["text"]),
        ("Поверхности", "#ffffff", COLORS["surface"], COLORS["text"]),
        ("Основной текст", "#1f2933", COLORS["text"], COLORS["surface"]),
        ("Приглушённый текст", "#6b7280", COLORS["muted"], COLORS["surface"]),
    ]

    swatch_size = Inches(1.5)
    start_x = Inches(1.0)
    start_y = Inches(1.6)
    gap = Inches(1.0)

    for i, (name, hex_code, color, text_color) in enumerate(palette):
        row = i // 3
        col = i % 3
        x = start_x + col * (swatch_size + gap)
        y = start_y + row * (swatch_size + Inches(0.9))

        # Цветной квадрат
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x), int(y),
            int(swatch_size), int(swatch_size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = COLORS["border"]
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.1

        # HEX-код внутри
        add_title(slide, hex_code,
                  x, y + swatch_size / 2 - Inches(0.15),
                  swatch_size, Inches(0.35),
                  font_size=18, bold=True, color=text_color, align=PP_ALIGN.CENTER)

        # Название под квадратом
        add_body(slide, name,
                 x, y + swatch_size + Inches(0.1),
                 swatch_size, Inches(0.4),
                 font_size=14, color=COLORS["text"], align=PP_ALIGN.CENTER)

    add_body(slide, "Зелёный символизирует природу и учёбу, красный — историю и наследие, кремовый фон создаёт тёплую патриотичную атмосферу.",
             MARGIN, Inches(6.2),
             Inches(12), Inches(0.8),
             font_size=16, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    # ============================================================
    # Слайд 10: База данных
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["primary"], top=True)

    add_title(slide, "База данных SQLite",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    tables = [
        "users — аккаунты и баллы",
        "articles — статьи и вопросы (JSON)",
        "results — результаты викторин",
        "messages — личные сообщения",
        "duels — PvP-дуэли",
        "team_battles — групповые комнаты",
        "team_battle_participants — участники битв",
    ]

    for idx, txt in enumerate(tables):
        y = Inches(1.6) + idx * Inches(0.7)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(MARGIN + Inches(0.2)), int(y),
            int(Inches(6.5)), int(Inches(0.55))
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["surface"]
        shape.line.color.rgb = COLORS["primary"]
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.15
        add_body(slide, txt,
                 MARGIN + Inches(0.45), y + Inches(0.08),
                 Inches(6.0), Inches(0.4),
                 font_size=18, color=COLORS["text"])

    # Schema diagram text
    add_title(slide, "Связи",
              Inches(8.0), Inches(1.6),
              Inches(4.5), Inches(0.5),
              font_size=22, bold=True, color=COLORS["primary"])

    relations = [
        "results.user_id → users.id",
        "results.article_id → articles.id",
        "messages.sender_id → users.id",
        "messages.receiver_id → users.id",
        "duels.challenger_id → users.id",
        "team_battles.creator_id → users.id",
        "team_battle_participants.battle_id → team_battles.id",
    ]
    for idx, txt in enumerate(relations):
        add_body(slide, "• " + txt,
                 Inches(8.0), Inches(2.1) + idx * Inches(0.45),
                 Inches(4.5), Inches(0.4),
                 font_size=14, color=COLORS["text"])

    # ============================================================
    # Слайд 11: API
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["bg"])
    add_decorative_strip(slide, COLORS["accent"], top=True)

    add_title(slide, "REST API",
              MARGIN, Inches(0.5),
              Inches(12), Inches(0.8),
              font_size=36, bold=True, color=COLORS["primary"])

    api_sections = [
        ("Публичные", [
            "POST /api/register",
            "POST /api/login",
            "GET /api/articles",
            "GET /api/articles/:id",
            "GET /api/leaderboard",
        ], COLORS["primary"]),
        ("Авторизованные", [
            "GET /api/me",
            "POST /api/results",
            "GET /api/chat/messages",
            "POST /api/duels",
            "POST /api/team-battles",
        ], COLORS["accent"]),
        ("Админ", [
            "POST /api/admin/articles",
            "PUT /api/admin/articles/:id",
            "DELETE /api/admin/articles/:id",
            "GET /api/admin/users",
            "DELETE /api/admin/users/:id",
        ], COLORS["text"]),
    ]

    col_width = Inches(4.0)
    for col_idx, (title, endpoints, color) in enumerate(api_sections):
        x = MARGIN + col_idx * (col_width + Inches(0.15))

        add_title(slide, title,
                  x, Inches(1.5),
                  col_width, Inches(0.5),
                  font_size=20, bold=True, color=color)

        for row_idx, endpoint in enumerate(endpoints):
            y = Inches(2.1) + row_idx * Inches(0.7)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(x), int(y),
                int(col_width), int(Inches(0.6))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLORS["surface"]
            shape.line.color.rgb = color
            shape.line.width = Pt(1)
            shape.adjustments[0] = 0.1
            add_body(slide, endpoint,
                     x + Inches(0.15), y + Inches(0.1),
                     col_width - Inches(0.3), Inches(0.4),
                     font_size=14, color=COLORS["text"])

    # ============================================================
    # Слайд 12: Спасибо
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, COLORS["primary"])

    add_title(slide, "Спасибо за внимание!",
              MARGIN, Inches(2.2),
              Inches(12), Inches(1),
              font_size=48, bold=True, color=COLORS["surface"], align=PP_ALIGN.CENTER)

    add_body(slide, "Спадчына — узнай свою страну интерактивно",
             MARGIN, Inches(3.4),
             Inches(12), Inches(0.8),
             font_size=24, color=COLORS["surface"], align=PP_ALIGN.CENTER)

    add_body(slide, "CultCode / Repablik\nФронтенд: http://localhost:5173\nБэкенд: http://127.0.0.1:8081",
             MARGIN, Inches(4.5),
             Inches(12), Inches(1.5),
             font_size=16, color=COLORS["surface"], align=PP_ALIGN.CENTER)

    # Сохранение
    output_path = "/Users/sergej/Desktop/codex-projects/repablik/Спадчына_презентация.pptx"
    prs.save(output_path)
    print(f"Презентация сохранена: {output_path}")


if __name__ == "__main__":
    create_presentation()
